from datetime import datetime, timezone
from typing import Annotated

from fastapi import APIRouter, File, Form, UploadFile, HTTPException
from pptx import Presentation
from docx import Document
from io import BytesIO
from urllib.parse import urlparse, parse_qs
import re
import httpx

from db import documents_collection

from config import settings

router = APIRouter()


@router.post("/upload")
async def upload(
    files: Annotated[list[UploadFile] | None, File()] = None,
    url: Annotated[str | None, Form()] = None,
):
    files_given = files is not None and len(files) > 0
    url_given = url is not None and url.strip() != ""

    if not files_given and not url_given:
        raise HTTPException(
            status_code=422,
            detail="Provide at least one file or a YouTube URL.",
        )

    document_id = None

    if url_given:
        parsed = urlparse(url)
        if parsed.netloc not in {"www.youtube.com", "youtube.com", "youtu.be"}:
            raise HTTPException(status_code=422, detail="Invalid YouTube URL.")

        text = extract_youtube_text(url)

        result = documents_collection.insert_one(
            {
                "name": url,
                "extracted_text": text,
                "file_type": "youtube",
                "upload_time": datetime.now(timezone.utc),
            }
        )
        document_id = str(result.inserted_id)

    if files_given:
        for f in files:
            if not f.filename.lower().endswith((".docx", ".pptx")):
                raise HTTPException(
                    status_code=415,
                    detail=f"Unsupported file type: {f.filename}",
                )

            if f.filename.lower().endswith(".docx"):
                text = extract_docx_text(f.file)
            else:
                text = extract_pptx_text(f.file)

            ext = f.filename.rsplit(".", 1)[-1]
            ext = re.sub(r"\s+", "", ext, flags=re.UNICODE).lower()

            result = documents_collection.insert_one(
                {
                    "name": f.filename,
                    "extracted_text": text,
                    "file_type": ext,
                    "upload_time": datetime.now(timezone.utc),
                }
            )
            document_id = str(result.inserted_id)

    return {"message": "OK", "document_id": document_id}


def extract_youtube_text(url: str) -> str:
    parsed_url = urlparse(url)

    video_id = None
    if parsed_url.netloc in {"www.youtube.com", "youtube.com"}:
        video_id = parse_qs(parsed_url.query).get("v", [None])[0]
    elif parsed_url.netloc == "youtu.be":
        video_id = parsed_url.path.lstrip("/")

    if not video_id:
        raise HTTPException(status_code=422, detail="Invalid YouTube URL.")

    api_url = f"https://{settings.rapidapi_host}/api/transcript"
    headers = {
        "x-rapidapi-key": settings.rapidapi_key,
        "x-rapidapi-host": settings.rapidapi_host,
    }
    params = {"videoId": video_id}

    try:
        r = httpx.get(api_url, headers=headers, params=params, timeout=20.0)
    except httpx.RequestError as e:
        raise HTTPException(status_code=503, detail=f"Transcript service unavailable: {type(e).__name__}")

    if r.status_code != 200:
        try:
            detail = r.json()
        except Exception:
            detail = r.text
        raise HTTPException(status_code=422, detail={"error": "Could not fetch transcript", "upstream": detail})

    data = r.json()

    items = None
    if isinstance(data, dict):
        if "transcript" in data and isinstance(data["transcript"], list):
            items = data["transcript"]
        elif "segments" in data and isinstance(data["segments"], list):
            items = data["segments"]
        elif "data" in data and isinstance(data["data"], list):
            items = data["data"]

    if not items:
        raise HTTPException(status_code=422, detail={"error": "Unexpected transcript response shape", "raw": data})

    text_parts: list[str] = []
    for item in items:
        if isinstance(item, dict) and "text" in item and isinstance(item["text"], str):
            text_parts.append(item["text"])

    return " ".join(text_parts).strip()


def extract_pptx_text(file_obj) -> str:
    prs = Presentation(BytesIO(file_obj.read()))
    text_runs: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        text_runs.append(run.text)

    return " ".join(text_runs)


def extract_docx_text(file_obj) -> str:
    doc = Document(BytesIO(file_obj.read()))
    text_parts: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    text_parts.append(cell.text)

    return " ".join(text_parts)