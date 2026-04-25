import io
import pytest
from unittest.mock import patch, MagicMock

from tests.conftest import AUTH, documents_collection


def _docx_bytes() -> bytes:
    """Minimal valid .docx file created with python-docx."""
    from docx import Document
    buf = io.BytesIO()
    doc = Document()
    doc.add_paragraph("SQL injection overview.")
    doc.save(buf)
    return buf.getvalue()


def _pptx_bytes() -> bytes:
    """Minimal valid .pptx file created with python-pptx."""
    from pptx import Presentation
    buf = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "SQL Injection"
    prs.save(buf)
    return buf.getvalue()


class TestUpload:
    def test_upload_docx_returns_document_id(self, client):
        documents_collection.insert_one.return_value.inserted_id = "507f1f77bcf86cd799439011"

        response = client.post(
            "/upload",
            files={"files": ("notes.docx", io.BytesIO(_docx_bytes()), "application/octet-stream")},
            headers=AUTH,
        )

        assert response.status_code == 200
        data = response.json()
        assert "document_id" in data
        assert data["document_id"] == "507f1f77bcf86cd799439011"

    def test_upload_pptx_returns_document_id(self, client):
        documents_collection.insert_one.return_value.inserted_id = "507f1f77bcf86cd799439012"

        response = client.post(
            "/upload",
            files={"files": ("slides.pptx", io.BytesIO(_pptx_bytes()), "application/octet-stream")},
            headers=AUTH,
        )

        assert response.status_code == 200
        assert response.json()["document_id"] == "507f1f77bcf86cd799439012"

    def test_upload_docx_saves_extracted_text_to_db(self, client):
        documents_collection.insert_one.return_value.inserted_id = "abc123"

        client.post(
            "/upload",
            files={"files": ("notes.docx", io.BytesIO(_docx_bytes()), "application/octet-stream")},
            headers=AUTH,
        )

        call_doc = documents_collection.insert_one.call_args[0][0]
        assert "extracted_text" in call_doc
        assert "SQL injection overview." in call_doc["extracted_text"]

    def test_upload_docx_stores_correct_file_type(self, client):
        documents_collection.insert_one.return_value.inserted_id = "abc"

        client.post(
            "/upload",
            files={"files": ("notes.docx", io.BytesIO(_docx_bytes()), "application/octet-stream")},
            headers=AUTH,
        )

        call_doc = documents_collection.insert_one.call_args[0][0]
        assert call_doc["file_type"] == "docx"

    def test_upload_pptx_stores_correct_file_type(self, client):
        documents_collection.insert_one.return_value.inserted_id = "abc"

        client.post(
            "/upload",
            files={"files": ("slides.pptx", io.BytesIO(_pptx_bytes()), "application/octet-stream")},
            headers=AUTH,
        )

        call_doc = documents_collection.insert_one.call_args[0][0]
        assert call_doc["file_type"] == "pptx"

    def test_upload_no_input_returns_422(self, client):
        response = client.post("/upload", headers=AUTH)
        assert response.status_code == 422

    def test_upload_unsupported_file_type_returns_415(self, client):
        response = client.post(
            "/upload",
            files={"files": ("report.pdf", io.BytesIO(b"%PDF-1.4"), "application/pdf")},
            headers=AUTH,
        )
        assert response.status_code == 415

    def test_upload_invalid_youtube_url_returns_422(self, client):
        response = client.post(
            "/upload",
            data={"url": "https://www.notayoutubesite.com/watch?v=abc"},
            headers=AUTH,
        )
        assert response.status_code == 422

    def test_upload_youtube_url_calls_transcript_api(self, client):
        documents_collection.insert_one.return_value.inserted_id = "yt123"

        transcript_response = MagicMock()
        transcript_response.status_code = 200
        transcript_response.json.return_value = {
            "transcript": [{"text": "Hello"}, {"text": "world"}]
        }

        with patch("routes.upload.httpx.get", return_value=transcript_response):
            response = client.post(
                "/upload",
                data={"url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ"},
                headers=AUTH,
            )

        assert response.status_code == 200
        assert response.json()["document_id"] == "yt123"

    def test_upload_youtube_saves_transcript_text(self, client):
        documents_collection.insert_one.return_value.inserted_id = "yt456"

        transcript_response = MagicMock()
        transcript_response.status_code = 200
        transcript_response.json.return_value = {
            "transcript": [{"text": "Photosynthesis"}, {"text": "is cool"}]
        }

        with patch("routes.upload.httpx.get", return_value=transcript_response):
            client.post(
                "/upload",
                data={"url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ"},
                headers=AUTH,
            )

        call_doc = documents_collection.insert_one.call_args[0][0]
        assert "Photosynthesis" in call_doc["extracted_text"]

    def test_upload_requires_auth(self, client):
        response = client.post(
            "/upload",
            files={"files": ("notes.docx", io.BytesIO(_docx_bytes()), "application/octet-stream")},
        )
        assert response.status_code in (401, 422)

    def test_youtube_transcript_api_unavailable_returns_503(self, client):
        import httpx
        with patch("routes.upload.httpx.get", side_effect=httpx.RequestError("connection refused")):
            response = client.post(
                "/upload",
                data={"url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ"},
                headers=AUTH,
            )
        assert response.status_code == 503

    def test_youtube_transcript_api_non_200_returns_422(self, client):
        error_response = MagicMock()
        error_response.status_code = 500
        error_response.json.return_value = {"message": "Internal Server Error"}
        with patch("routes.upload.httpx.get", return_value=error_response):
            response = client.post(
                "/upload",
                data={"url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ"},
                headers=AUTH,
            )
        assert response.status_code == 422
